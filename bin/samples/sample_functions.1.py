#!/usr/bin/env python3
# -*- coding: utf-8 -*-

"""
Style:
    Google Python Style Guide:
    http://google.github.io/styleguide/pyguide.html
    @name           sample_functions
    @version        1.0.0
    @author-name    Wayne Schmidt
    @author-email   wschmidt@sumologic.com
    @license-name   APACHE 2.0
    @license-url    http://www.apache.org/licenses/LICENSE-2.0
"""

__version__ = 1.0
__author__ = "Wayne Schmidt (wschmidt@sumologic.com)"

import argparse
import glob
import os
import sys
import re
import datetime
from itertools import chain
import pandas as pd

from pptx import Presentation
from pptx.util import Inches
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

### from pptx.enum.dml import MSO_THEME_COLOR
### from pptx.enum.text import MSO_AUTO_SIZE

sys.dont_write_bytecode = 1

NOW = datetime.datetime.now()
TSTAMP = NOW.strftime("%B %-d, %Y")

STATUS = {}
STATUS['COMPLETED'] = 0x00, 0xFF, 0x00
STATUS['IN_PROGRESS'] = 0xFF, 0x99, 0x00
STATUS['DELAYED'] = 0xFF, 0x00, 0x00
STATUS['NEW'] = 0x99, 0xCC, 0xFF

MILESTONE = {}
MILESTONE['VALUE'] = 0x00, 0xFF, 0x00
MILESTONE['CONTENT'] = 0x00, 0xFF, 0x00
MILESTONE['ONBOARD'] = 0xFF, 0xFF, 0x00
MILESTONE['REVIEW'] = 0xFF, 0xFF, 0x00
MILESTONE['CHECKUP'] = 0xFF, 0x99, 0x00
MILESTONE['CSP'] = 0xFF, 0x99, 0x00

def add_front_logo(path, img, left, top, width):
    """
    This adds a front logo to the 1st slide defined in the template
    """
    logo_x = Inches(left)
    logo_y = Inches(top)
    logo_w = Inches(width)

    logo_slide = PRESENTATION.slides[0]
    _logo_pic = logo_slide.shapes.add_picture(os.path.join(path, img), logo_x, logo_y, logo_w)

def add_front_title(text, left, top, width, height):
    """
    This adds a front title to the 1st slide defined in the template
    """
    text_x = Inches(left)
    text_y = Inches(top)
    text_w = Inches(width)
    text_h = Inches(height)

    text_slide = PRESENTATION.slides[0]
    tx_box = text_slide.shapes.add_textbox(text_x, text_y, text_w, text_h)
    t_f = tx_box.text_frame
    t_para = t_f.add_paragraph()
    t_para.text = text
    t_para.font.size = Pt(36)
    t_para.font.name = 'Calibri'
    t_para.alignment = PP_ALIGN.RIGHT

def add_support_table(cfgpath, userfile, left, top, width, height):
    """
    This will add a support table based on a CSV file.
    Eventually this will be replaced by a feed from Salesforce
    """
    table_x = Inches(left)
    table_y = Inches(top)
    table_cx = Inches(width)
    table_cy = Inches(height)

    dataframe = pd.read_csv(os.path.abspath(os.path.join(cfgpath, userfile)))
    (rows, columns) = dataframe.shape
    rows = rows + 1

    shape = SLIDE.shapes.add_table(rows, columns, table_x, table_y, table_cx, table_cy)
    table = shape.table
    set_table_contents(dataframe, table)
    add_table_markup(table)

def set_table_contents(dataframe, table):
    """
    This sets the color and other attributes
    """
    i_c = 0
    for col, row in dataframe.iteritems():
        i_r = 0
        cell = table.cell(i_r, i_c)
        cell.text = col
        i_r = i_r + 1
        for data in row:
            cell = table.cell(i_r, i_c)
            cell.text = str(data)
            if cell.text == "nan":
                cell.text = ""
            if cell.text in STATUS:
                cell.fill.solid()
                (r_1, c_1, b_1) = STATUS[cell.text]
                cell.fill.fore_color.rgb = RGBColor(r_1, c_1, b_1)
            if cell.text in MILESTONE:
                cell.fill.solid()
                (r_1, c_1, b_1) = MILESTONE[cell.text]
                cell.fill.fore_color.rgb = RGBColor(r_1, c_1, b_1)
            i_r = i_r + 1
        i_c = i_c + 1

def iter_cells(table):
    """
    calculate the cell for a given table
    """
    for row in table.rows:
        for cell in row.cells:
            yield cell

def add_table_markup(table):
    """
    Adjusts font size and adds hyperlinks
    """
    for cell in iter_cells(table):
        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(16)
                run.font.name = 'Calibri'
                if re.search(r'https://*', cell.text):
                    ticket = str(cell.text.split('/')[-1])
                    run.text = ticket
                    hlink = run.hyperlink
                    hlink.address = cell.text
                if re.search(r'@', cell.text):
                    email = str(cell.text.split('@')[0])
                    domain = str(cell.text.split('@')[1])
                    run.text = email
                    hlink = run.hyperlink
                    hlink.address = 'mailto:' + email + '@' + domain

def add_logo(path, img, left, top, width):
    """
    Add a logo to a specific slide
    """
    logo_x = Inches(left)
    logo_y = Inches(top)
    logo_w = Inches(width)
    _logo_pic = SLIDE.shapes.add_picture(os.path.join(path, img), logo_x, logo_y, logo_w)


def add_title(text, left, top, width, height):
    """
    Add a title to a given slide
    """
    logo_x = Inches(left)
    logo_y = Inches(top)
    logo_w = Inches(width)
    logo_h = Inches(height)

    tx_box = SLIDE.shapes.add_textbox(logo_x, logo_y, logo_w, logo_h)
    t_f = tx_box.text_frame
    t_para = t_f.add_paragraph()
    t_para.text = text
    t_para.font.size = Pt(36)
    t_para.font.name = 'Calibri'
    t_para.alignment = PP_ALIGN.LEFT

def cleanup_slides(prs):
    """
    Remove Slides in Template. Workaround for duplication or renumbering
    """
    prs_slides = prs.slides
    num_pages = len(prs_slides)

    top_range = range(num_pages - 1, 3, -1)
    all_range = chain(top_range)

    for _i in all_range:
        r_id = prs.slides._sldIdLst[-1].rId
        prs.part.drop_rel(r_id)
        del prs.slides._sldIdLst[-1]

def validate_template(tmplpath):
    TMPLPATH = ARGS.template
    if os.path.isdir(TMPLPATH):
        TMPLNAME = glob.glob(os.path.join(os.path.abspath(TMPLPATH), '*.pptx'))[0]
    else:
        print('Path not accessible:: {} '.format(TMPLPATH))
        sys.exit()
