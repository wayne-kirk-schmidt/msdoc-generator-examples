#!/usr/bin/env python3
# -*- coding: utf-8 -*-

"""
Style:
    Google Python Style Guide:
    http://google.github.io/styleguide/pyguide.html
    @name           sample_functions
    @version        1.0.0
    @author-name    Wayne Schmidt
    @author-email   wschmidt@sumologic.com
    @license-name   APACHE 2.0
    @license-url    http://www.apache.org/licenses/LICENSE-2.0
"""

__version__ = 1.0
__author__ = "Wayne Schmidt (wschmidt@sumologic.com)"

import argparse
import glob
import os
import sys
import re
import datetime
from itertools import chain
import pandas as pd

from pptx import Presentation
from pptx.util import Inches
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

### from pptx.enum.dml import MSO_THEME_COLOR
### from pptx.enum.text import MSO_AUTO_SIZE

sys.dont_write_bytecode = 1

my_file_name = sys.argv[1]
my_target = int(sys.argv[2])
my_presentation = Presentation(my_file_name)

def get_slide_count(presentation):
    slidecount = 0
    for slide in presentation.slides:
        slidecount += 1
    return slidecount

def delete_slide(presentation, slide):
    id_dict = { slide.id: [i, slide.rId] for i,slide in enumerate(presentation.slides._sldIdLst) }
    slide_id = slide.slide_id
    presentation.part.drop_rel(id_dict[slide_id][1])
    del presentation.slides._sldIdLst[id_dict[slide_id][0]]

def get_single_slide(presentation, target):
    for idx, slide in enumerate(presentation.slides):
        if idx < my_target:
            delete_slide(presentation, slide)
        elif (idx > my_target):
            delete_slide(presentation, slide)
    presentation.save("sample-" + str(target + 1) + ".pptx")

slidecount = get_slide_count(my_presentation)

for counter in range(slidecount):
    my_backup= Presentation(my_file_name)
    get_single_slide(my_backup, counter)
    my_backup = None

